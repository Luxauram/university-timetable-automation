"""
converter.py
============
PDF-to-PowerPoint conversion layer.

Uses pypdfium2 (Python bindings for Google's PDFium engine) to render each
PDF page as a high-resolution in-memory PNG image, then assembles those images
into a .pptx file via python-pptx.

No system-level dependencies are required (no Poppler, Ghostscript or
ImageMagick). pypdfium2 ships its own pre-compiled PDFium binary and works
out of the box on macOS 12+, Windows, and Linux.

License: MIT
"""

import io

import pypdfium2 as pdfium
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from config import PPTX_DPI, PPTX_SLIDE_W_IN, PPTX_SLIDE_H_IN


def pdf_to_pptx(pdf_path: str, pptx_path: str) -> bool:
    """
    Converts a PDF file into a PowerPoint presentation.

    Each page of the source PDF is rendered as a full-slide image and placed
    on a blank slide. The resulting .pptx file uses A4 landscape dimensions
    as defined in config.py.

    @param pdf_path:  Absolute or relative path to the source PDF file.
    @param pptx_path: Destination path for the generated .pptx file.
    @return:          True if the conversion completed successfully,
                      False if any error occurred (error is printed to stdout).
    """
    try:
        page_images = _render_pdf_pages(pdf_path)
        if not page_images:
            return False
        _build_pptx(page_images, pptx_path)
        return True
    except Exception as exc:
        print(f"[converter] Conversion failed for '{pdf_path}': {exc}")
        return False


def _render_pdf_pages(pdf_path: str) -> list[io.BytesIO]:
    """
    Opens a PDF with PDFium and renders every page to an in-memory PNG buffer.

    The rendering scale is derived from PPTX_DPI (config.py) using the
    formula ``scale = DPI / 72``, since PDFium's internal coordinate system
    is based on 72 points per inch.

    Pages are closed immediately after rendering to release native memory.
    The document is always closed in the finally block, even on error.

    @param pdf_path: Path to the PDF file to render.
    @return:         Ordered list of BytesIO buffers, one PNG image per page.
    @raises pdfium.PdfiumError: If the file cannot be opened or is corrupted.
    """
    scale = PPTX_DPI / 72
    doc = pdfium.PdfDocument(pdf_path)
    images: list[io.BytesIO] = []

    try:
        for index in range(len(doc)):
            page = doc[index]
            bitmap = page.render(scale=scale, rotation=0)
            pil_image = bitmap.to_pil()
            page.close()

            buffer = io.BytesIO()
            pil_image.save(buffer, format="PNG")
            buffer.seek(0)
            images.append(buffer)
    finally:
        doc.close()

    return images


def _build_pptx(page_images: list[io.BytesIO], pptx_path: str) -> None:
    """
    Assembles a list of in-memory PNG images into a .pptx file.

    Slide dimensions are read from config.py (PPTX_SLIDE_W_IN,
    PPTX_SLIDE_H_IN). Each image is placed on a separate blank slide and
    centred using proportional scaling (aspect ratio is always preserved).

    @param page_images: Ordered list of BytesIO PNG buffers, one per slide.
    @param pptx_path:   Destination file path for the saved presentation.
    """
    presentation = Presentation()
    presentation.slide_width = Inches(PPTX_SLIDE_W_IN)
    presentation.slide_height = Inches(PPTX_SLIDE_H_IN)

    blank_layout = presentation.slide_layouts[6]

    for image_buffer in page_images:
        slide = presentation.slides.add_slide(blank_layout)
        _place_image_centered(
            slide, image_buffer, presentation.slide_width, presentation.slide_height
        )

    presentation.save(pptx_path)


def _place_image_centered(slide, image_buffer: io.BytesIO, slide_w: int, slide_h: int) -> None:
    """
    Places a PNG image onto a slide, scaled to fit and centred both
    horizontally and vertically without distortion.

    The image buffer seek position is reset before reading and again before
    being passed to python-pptx, so the same buffer can be safely reused
    by the caller.

    @param slide:        The python-pptx Slide object to draw on.
    @param image_buffer: BytesIO buffer containing a PNG-encoded image.
    @param slide_w:      Slide width in EMU (English Metric Units).
    @param slide_h:      Slide height in EMU.
    """
    image_buffer.seek(0)
    with Image.open(image_buffer) as im:
        img_w, img_h = im.size
    image_buffer.seek(0)

    ratio = min(slide_w / img_w, slide_h / img_h)
    new_w = int(img_w * ratio)
    new_h = int(img_h * ratio)
    left = (slide_w - new_w) // 2
    top = (slide_h - new_h) // 2

    slide.shapes.add_picture(image_buffer, left, top, new_w, new_h)
